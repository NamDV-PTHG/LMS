"""
Document parser — supports PDF, DOCX, PPTX.
Returns chunked text lists (~500 tokens each).
"""
import re
import io


# ── Chunker ───────────────────────────────────────────────────

def chunk_text(text: str, max_tokens: int = 500) -> list[str]:
    """
    Split text into chunks not exceeding max_tokens (approx. 4 chars/token).
    Splits on sentence boundaries where possible.
    """
    max_chars = max_tokens * 4
    # Split on sentence-ending punctuation
    sentences = re.split(r'(?<=[.!?\n])\s+', text.strip())

    chunks: list[str] = []
    current = ""

    for sentence in sentences:
        if not sentence.strip():
            continue
        if len(current) + len(sentence) + 1 <= max_chars:
            current = (current + " " + sentence).strip()
        else:
            if current:
                chunks.append(current)
            # If single sentence exceeds max, force split
            if len(sentence) > max_chars:
                for i in range(0, len(sentence), max_chars):
                    chunks.append(sentence[i : i + max_chars])
                current = ""
            else:
                current = sentence

    if current:
        chunks.append(current)

    return [c for c in chunks if len(c.strip()) > 50]  # drop tiny fragments


# ── PDF ───────────────────────────────────────────────────────

async def parse_pdf(file_bytes: bytes) -> list[str]:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("pymupdf is required: pip install pymupdf")

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    chunks: list[str] = []

    for page in doc:
        text = page.get_text("text")
        if text.strip():
            chunks.extend(chunk_text(text))

    doc.close()
    return chunks


# ── DOCX ──────────────────────────────────────────────────────

async def parse_docx(file_bytes: bytes) -> list[str]:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx is required: pip install python-docx")

    doc = Document(io.BytesIO(file_bytes))
    full_text = "\n".join(
        p.text for p in doc.paragraphs if p.text.strip()
    )
    return chunk_text(full_text)


# ── PPTX ──────────────────────────────────────────────────────

async def parse_pptx(file_bytes: bytes) -> list[str]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required: pip install python-pptx")

    prs = Presentation(io.BytesIO(file_bytes))
    chunks: list[str] = []

    for slide in prs.slides:
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_parts.append(shape.text.strip())
        if slide_text_parts:
            slide_text = "\n".join(slide_text_parts)
            chunks.extend(chunk_text(slide_text))

    return chunks


# ── Dispatcher ────────────────────────────────────────────────

async def parse_document(file_bytes: bytes, mime_type: str) -> list[str]:
    if mime_type == "application/pdf":
        return await parse_pdf(file_bytes)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return await parse_docx(file_bytes)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return await parse_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {mime_type}")
